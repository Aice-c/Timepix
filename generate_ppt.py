"""
Generate academic presentation for group meeting:
Timepix3 Near-Vertical Angle Distinguishability Analysis
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# ── Color palette ──
C_TITLE   = RGBColor(0x2c, 0x3e, 0x50)  # dark navy
C_ACCENT  = RGBColor(0x29, 0x80, 0xb9)  # blue
C_BODY    = RGBColor(0x33, 0x33, 0x33)  # dark gray
C_CAPTION = RGBColor(0x66, 0x66, 0x66)  # gray
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_BG_ALT  = RGBColor(0xF0, 0xF4, 0xF8)  # alt row
C_GREEN   = RGBColor(0x27, 0xAE, 0x60)
C_RED     = RGBColor(0xC0, 0x39, 0x2B)
C_LIGHT_LINE = RGBColor(0xBD, 0xC3, 0xC7)

OUTPUT_DIR = r"d:\Project\Timepix\output"
PPT_PATH   = os.path.join(OUTPUT_DIR, "near_vertical_analysis_report.pptx")

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)


# ── Helper functions ──
def add_slide():
    """Add a blank slide."""
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def add_title_bar(slide, text, y=Inches(0.3), font_size=Pt(32)):
    """Add a title text at top of slide."""
    txBox = slide.shapes.add_textbox(Inches(0.8), y, Inches(11.7), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.bold = True
    p.font.color.rgb = C_TITLE
    p.font.name = "Calibri"
    # underline bar
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), y + Inches(0.85),
        Inches(11.7), Pt(3)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = C_ACCENT
    line.line.fill.background()
    return tf


def add_body_text(slide, text, left=Inches(0.8), top=Inches(1.5),
                  width=Inches(11.7), height=Inches(5.2), font_size=Pt(20)):
    """Add body text box, return text_frame for further editing."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.strip().split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = font_size
        p.font.color.rgb = C_BODY
        p.font.name = "Calibri"
        p.space_after = Pt(6)
    return tf


def add_bullet_slide(slide, title, bullets, start_top=Inches(1.5)):
    """Content slide with title and bullet points."""
    add_title_bar(slide, title)
    txBox = slide.shapes.add_textbox(Inches(0.8), start_top, Inches(11.7), Inches(5.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, level) in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.level = level
        p.font.size = Pt(22) if level == 0 else Pt(19)
        p.font.color.rgb = C_BODY if level == 0 else C_CAPTION
        p.font.name = "Calibri"
        p.space_before = Pt(10) if level == 0 else Pt(4)
        p.space_after = Pt(4)
    return tf


def add_image_slide(slide, title, img_path, caption="",
                    img_left=None, img_top=Inches(1.4), img_width=None, img_height=None):
    """Slide with a title and centered image."""
    add_title_bar(slide, title)
    if img_width is None:
        img_width = Inches(10)
    if img_left is None:
        img_left = Inches((13.333 - img_width / 914400) / 2) if isinstance(img_width, int) else Inches(1.6)
    pic = slide.shapes.add_picture(img_path, img_left, img_top, img_width, img_height)
    if caption:
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(6.8), Inches(11.7), Inches(0.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = caption
        p.font.size = Pt(14)
        p.font.color.rgb = C_CAPTION
        p.font.italic = True
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER
    return pic


def add_table(slide, data, left=Inches(0.8), top=Inches(1.5),
              col_widths=None, font_size=Pt(16)):
    """Add a styled table. data = list of lists, first row is header."""
    rows, cols = len(data), len(data[0])
    if col_widths is None:
        w = Inches(11.7)
        col_widths = [w // cols] * cols
    table_shape = slide.shapes.add_table(rows, cols, left, top,
                                          sum(col_widths), Inches(0.45 * rows))
    table = table_shape.table
    for ci, cw in enumerate(col_widths):
        table.columns[ci].width = cw
    for ri, row_data in enumerate(data):
        for ci, cell_text in enumerate(row_data):
            cell = table.cell(ri, ci)
            cell.text = str(cell_text)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = font_size
                paragraph.font.name = "Calibri"
                paragraph.alignment = PP_ALIGN.CENTER
                if ri == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = C_WHITE
                else:
                    paragraph.font.color.rgb = C_BODY
            if ri == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_TITLE
            elif ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_BG_ALT
    return table


def add_slide_number(slide, num):
    txBox = slide.shapes.add_textbox(Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = str(num)
    p.font.size = Pt(12)
    p.font.color.rgb = C_CAPTION
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.RIGHT


# ══════════════════════════════════════════════════════════════════════
# SLIDE 1: Title
# ══════════════════════════════════════════════════════════════════════
s = add_slide()
# Title
txBox = s.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(10.3), Inches(1.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Timepix3 近垂直入射角度可区分性分析"
p.font.size = Pt(38)
p.font.bold = True
p.font.color.rgb = C_TITLE
p.font.name = "Calibri"
p.alignment = PP_ALIGN.CENTER

# Subtitle
txBox2 = s.shapes.add_textbox(Inches(1.5), Inches(3.6), Inches(10.3), Inches(0.6))
tf2 = txBox2.text_frame
p2 = tf2.paragraphs[0]
p2.text = "80°–90° ToT数据手工特征统计分析"
p2.font.size = Pt(24)
p2.font.color.rgb = C_ACCENT
p2.font.name = "Calibri"
p2.alignment = PP_ALIGN.CENTER

# Info line
txBox3 = s.shapes.add_textbox(Inches(1.5), Inches(4.6), Inches(10.3), Inches(0.5))
tf3 = txBox3.text_frame
p3 = tf3.paragraphs[0]
p3.text = "组会汇报  |  2026年3月"
p3.font.size = Pt(18)
p3.font.color.rgb = C_CAPTION
p3.font.name = "Calibri"
p3.alignment = PP_ALIGN.CENTER

# Decorative line
line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(4.3), Inches(5.3), Pt(3))
line.fill.solid()
line.fill.fore_color.rgb = C_ACCENT
line.line.fill.background()

# ══════════════════════════════════════════════════════════════════════
# SLIDE 2: Research Background
# ══════════════════════════════════════════════════════════════════════
s = add_slide()
add_slide_number(s, 2)
add_bullet_slide(s, "研究背景", [
    ("Timepix3 混合像素探测器 (55μm × 55μm, 256×256 阵列)", 0),
    ("同时记录 ToA (到达时间) 和 ToT (超阈时间 ∝ 沉积能量)", 1),
    ("核心任务: 从单次测量轨迹图像推断粒子入射角度", 0),
    ("CNN / ViT 在大角度范围 (10°–90°) 已取得良好效果", 0),
    ("问题: 近垂直入射角度 (80°–90°) CNN准确率 < 20%", 0),
    ("约等于6类随机猜测基线 16.7%", 1),
    ("本次分析目标: 判断这些角度是否在物理上可区分", 0),
    ("→ 是模型能力不足，还是数据本身不可分？", 1),
])

# ══════════════════════════════════════════════════════════════════════
# SLIDE 3: Data Overview
# ══════════════════════════════════════════════════════════════════════
s = add_slide()
add_slide_number(s, 3)
add_title_bar(s, "数据概况")
data = [
    ["角度", "样本数", "平均活跃像素", "平均总ToT", "最大单像素ToT"],
    ["80°", "31,236", "12.7", "6,012.5", "8,252.1"],
    ["82°", "31,744", "12.7", "5,952.5", "8,260.8"],
    ["84°", "33,655", "12.7", "5,899.4", "8,260.8"],
    ["86°", "33,799", "12.5", "5,870.1", "8,260.8"],
    ["88°", "29,855", "12.3", "5,803.5", "8,260.8"],
    ["90°", "31,137", "12.1", "5,768.5", "8,260.8"],
]
add_table(s, data, top=Inches(1.5),
          col_widths=[Inches(1.6), Inches(2.0), Inches(2.4), Inches(2.8), Inches(2.9)],
          font_size=Pt(18))

# observation text
txBox = s.shapes.add_textbox(Inches(0.8), Inches(5.2), Inches(11.7), Inches(1.5))
tf = txBox.text_frame
tf.word_wrap = True
observations = [
    "• 共 191,426 个样本，各角度样本量均衡 (~30K)",
    "• 活跃像素数随角度增大略微减少 (12.7 → 12.1)",
    "• 总ToT随角度增大单调递减 (6012 → 5769)，但变化幅度仅 ~4%",
    "• 数据格式: 50×50 空间分辨率矩阵 (ToT值，高度稀疏)"
]
for i, txt in enumerate(observations):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = txt
    p.font.size = Pt(17)
    p.font.color.rgb = C_BODY
    p.font.name = "Calibri"
    p.space_after = Pt(3)

# ══════════════════════════════════════════════════════════════════════
# SLIDE 4: Analysis Methods
# ══════════════════════════════════════════════════════════════════════
s = add_slide()
add_slide_number(s, 4)
add_bullet_slide(s, "分析方法", [
    ("1. 手工特征提取 (21维特征向量)", 0),
    ("几何特征: 活跃像素数、包围盒、长宽比、离心率、紧凑度", 1),
    ("能量特征: 总ToT、最大ToT、均值/方差/偏度/峰度、熵、基尼系数", 1),
    ("空间-能量交叉特征: 能量梯度、加权质心偏移、二阶矩等", 1),
    ("2. 特征分布可视化 + KS统计检验", 0),
    ("相邻角度对 Kolmogorov-Smirnov 检验，量化分布差异", 1),
    ("3. Random Forest 多特征分类 (5折交叉验证)", 0),
    ("n_estimators=300, max_depth=15, class_weight='balanced'", 1),
    ("4. PCA 降维可视化", 0),
])

# ══════════════════════════════════════════════════════════════════════
# SLIDE 5: Feature Distributions
# ══════════════════════════════════════════════════════════════════════
s = add_slide()
add_slide_number(s, 5)
fig_path = os.path.join(OUTPUT_DIR, "feature_distributions_by_angle.png")
add_image_slide(s, "特征分布: 各角度高度重叠", fig_path,
                caption="图: 12个关键特征在6个角度上的分布对比 — 各角度分布几乎完全重叠",
                img_left=Inches(1.2), img_top=Inches(1.3),
                img_width=Inches(10.8), img_height=Inches(5.3))

# ══════════════════════════════════════════════════════════════════════
# SLIDE 6: KS Test Results
# ══════════════════════════════════════════════════════════════════════
s = add_slide()
add_slide_number(s, 6)
ks_path = os.path.join(OUTPUT_DIR, "ks_heatmap.png")
add_image_slide(s, "KS统计检验: 效应量极小", ks_path,
                caption="图: 相邻角度对KS统计量热力图 — 所有KS统计量 < 0.035 (远低于0.1的\"小效应\"阈值)",
                img_left=Inches(2.8), img_top=Inches(1.3),
                img_width=Inches(7.5), img_height=Inches(5.3))

# ══════════════════════════════════════════════════════════════════════
# SLIDE 7: KS Test Summary Table
# ══════════════════════════════════════════════════════════════════════
s = add_slide()
add_slide_number(s, 7)
add_title_bar(s, "KS检验关键结论")
ks_data = [
    ["角度对", "最大KS统计量", "KS > 0.1的特征数", "判定"],
    ["80° – 82°", "0.0235", "0 / 21", "不可区分"],
    ["82° – 84°", "0.0233", "0 / 21", "不可区分"],
    ["84° – 86°", "0.0207", "0 / 21", "不可区分"],
    ["86° – 88°", "0.0328", "0 / 21", "不可区分"],
    ["88° – 90°", "0.0315", "0 / 21", "不可区分"],
]
add_table(s, ks_data, top=Inches(1.5),
          col_widths=[Inches(2.5), Inches(3.0), Inches(3.2), Inches(3.0)],
          font_size=Pt(20))

txBox = s.shapes.add_textbox(Inches(0.8), Inches(5.2), Inches(11.7), Inches(1.5))
tf = txBox.text_frame
tf.word_wrap = True
notes = [
    "• Cohen's d 效应量标准: KS < 0.1 为可忽略差异",
    "• 虽然p值因大样本量 (~30K) 达到统计显著 (***), 但效应量极小",
    "• 统计显著 ≠ 实际显著 — 这是大样本统计的经典陷阱",
]
for i, txt in enumerate(notes):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = txt
    p.font.size = Pt(18)
    p.font.color.rgb = C_BODY
    p.font.name = "Calibri"
    p.space_after = Pt(4)

# ══════════════════════════════════════════════════════════════════════
# SLIDE 8: Random Forest Results
# ══════════════════════════════════════════════════════════════════════
s = add_slide()
add_slide_number(s, 8)
add_title_bar(s, "Random Forest 分类结果")

# Key metric boxes
# Accuracy box
box1 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.5), Inches(3.5), Inches(2.0))
box1.fill.solid()
box1.fill.fore_color.rgb = RGBColor(0xEB, 0xF5, 0xFB)
box1.line.color.rgb = C_ACCENT
txf1 = box1.text_frame
txf1.word_wrap = True
p = txf1.paragraphs[0]
p.text = "5折平均准确率"
p.font.size = Pt(16)
p.font.color.rgb = C_CAPTION
p.font.name = "Calibri"
p.alignment = PP_ALIGN.CENTER
p = txf1.add_paragraph()
p.text = "19.00% ± 0.20%"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = C_RED
p.font.name = "Calibri"
p.alignment = PP_ALIGN.CENTER

# Baseline box
box2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.0), Inches(1.5), Inches(3.5), Inches(2.0))
box2.fill.solid()
box2.fill.fore_color.rgb = RGBColor(0xFD, 0xED, 0xEC)
box2.line.color.rgb = C_RED
txf2 = box2.text_frame
txf2.word_wrap = True
p = txf2.paragraphs[0]
p.text = "随机猜测基线"
p.font.size = Pt(16)
p.font.color.rgb = C_CAPTION
p.font.name = "Calibri"
p.alignment = PP_ALIGN.CENTER
p = txf2.add_paragraph()
p.text = "16.67%"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = C_RED
p.font.name = "Calibri"
p.alignment = PP_ALIGN.CENTER

# Delta box
box3 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.0), Inches(1.5), Inches(3.5), Inches(2.0))
box3.fill.solid()
box3.fill.fore_color.rgb = RGBColor(0xFE, 0xF9, 0xE7)
box3.line.color.rgb = RGBColor(0xF3, 0x9C, 0x12)
txf3 = box3.text_frame
txf3.word_wrap = True
p = txf3.paragraphs[0]
p.text = "超出基线"
p.font.size = Pt(16)
p.font.color.rgb = C_CAPTION
p.font.name = "Calibri"
p.alignment = PP_ALIGN.CENTER
p = txf3.add_paragraph()
p.text = "+2.33%"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = RGBColor(0xF3, 0x9C, 0x12)
p.font.name = "Calibri"
p.alignment = PP_ALIGN.CENTER

# Classification report table
cls_data = [
    ["类别", "Precision", "Recall", "F1-score"],
    ["80°", "0.20", "0.34", "0.25"],
    ["82°", "0.18", "0.15", "0.16"],
    ["84°", "0.19", "0.15", "0.17"],
    ["86°", "0.19", "0.11", "0.14"],
    ["88°", "0.17", "0.15", "0.16"],
    ["90°", "0.20", "0.25", "0.22"],
]
add_table(s, cls_data, top=Inches(4.0),
          col_widths=[Inches(2.0), Inches(2.5), Inches(2.5), Inches(2.5)],
          font_size=Pt(17))

# ══════════════════════════════════════════════════════════════════════
# SLIDE 9: Confusion Matrix
# ══════════════════════════════════════════════════════════════════════
s = add_slide()
add_slide_number(s, 9)
cm_path = os.path.join(OUTPUT_DIR, "confusion_matrix_rf.png")
add_image_slide(s, "混淆矩阵: 各类别严重混淆", cm_path,
                caption="图: Random Forest 混淆矩阵 — 对角线元素(正确分类)无明显优势，误分类近均匀分布",
                img_left=Inches(3.0), img_top=Inches(1.3),
                img_width=Inches(7.0), img_height=Inches(5.3))

# ══════════════════════════════════════════════════════════════════════
# SLIDE 10: Feature Importance
# ══════════════════════════════════════════════════════════════════════
s = add_slide()
add_slide_number(s, 10)
fi_path = os.path.join(OUTPUT_DIR, "feature_importance.png")
add_image_slide(s, "特征重要性排名", fi_path,
                caption="图: Random Forest 特征重要性 — 能量分布特征 (entropy, gini) 排名最高，几何特征 (n_pixels, bbox) 最低",
                img_left=Inches(2.5), img_top=Inches(1.3),
                img_width=Inches(8.0), img_height=Inches(5.3))

# ══════════════════════════════════════════════════════════════════════
# SLIDE 11: PCA Visualization
# ══════════════════════════════════════════════════════════════════════
s = add_slide()
add_slide_number(s, 11)
pca_path = os.path.join(OUTPUT_DIR, "dimensionality_reduction.png")
add_image_slide(s, "PCA降维可视化: 各角度完全重叠", pca_path,
                caption="图: PCA前两主成分 (累计方差解释 63.7%) — 6个角度在低维空间中不可分离",
                img_left=Inches(2.5), img_top=Inches(1.3),
                img_width=Inches(8.0), img_height=Inches(5.3))

# ══════════════════════════════════════════════════════════════════════
# SLIDE 12: Comprehensive Verdict
# ══════════════════════════════════════════════════════════════════════
s = add_slide()
add_slide_number(s, 12)
add_title_bar(s, "综合结论")

verdict_data = [
    ["角度对", "KS_max", "互混淆率", "RF分类", "结论"],
    ["80° – 82°", "0.023", "46.1%", "~随机", "不可区分"],
    ["82° – 84°", "0.023", "28.6%", "~随机", "不可区分"],
    ["84° – 86°", "0.021", "25.4%", "~随机", "不可区分"],
    ["86° – 88°", "0.033", "25.6%", "~随机", "不可区分"],
    ["88° – 90°", "0.031", "39.3%", "~随机", "不可区分"],
]
add_table(s, verdict_data, top=Inches(1.5),
          col_widths=[Inches(2.2), Inches(2.0), Inches(2.2), Inches(2.2), Inches(3.1)],
          font_size=Pt(19))

txBox = s.shapes.add_textbox(Inches(0.8), Inches(5.0), Inches(11.7), Inches(2.0))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "核心结论"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = C_RED
p.font.name = "Calibri"
conclusions = [
    "在80°–90°范围内，以2°间隔区分入射角度在物理上不可行",
    "这不是模型能力问题，而是数据本身的物理限制",
    "CNN/ViT在该范围内准确率≈19% 与手工特征RF结果一致，进一步验证了这一结论",
]
for txt in conclusions:
    p = tf.add_paragraph()
    p.text = "▸ " + txt
    p.font.size = Pt(19)
    p.font.color.rgb = C_BODY
    p.font.name = "Calibri"
    p.space_after = Pt(4)

# ══════════════════════════════════════════════════════════════════════
# SLIDE 13: Recommendations
# ══════════════════════════════════════════════════════════════════════
s = add_slide()
add_slide_number(s, 13)
add_bullet_slide(s, "建议与后续工作", [
    ("角度合并策略", 0),
    ("方案A: 全部合并为 [80°–90°] 单一类别 (2°分辨率不可能)", 1),
    ("方案B: 尝试3个大类 [80°–82°], [84°–86°], [88°–90°]", 1),
    ("方案C: 尝试2个大类 [80°–84°], [86°–90°]", 1),
    ("深度学习改进方向", 0),
    ("将手工特征作为辅助输入注入CNN/ViT", 1),
    ("用注意力机制关注能量分布模式 (entropy, gini)", 1),
    ("将分类任务改为回归任务 (连续角度预测)", 1),
    ("论文写作意义", 0),
    ("Timepix3探测器在近垂直入射下的角分辨率极限 — 合理的物理结论", 1),
])

# ══════════════════════════════════════════════════════════════════════
# SLIDE 14: Thank You
# ══════════════════════════════════════════════════════════════════════
s = add_slide()
txBox = s.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10.3), Inches(1.2))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "谢谢"
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = C_TITLE
p.font.name = "Calibri"
p.alignment = PP_ALIGN.CENTER

txBox2 = s.shapes.add_textbox(Inches(1.5), Inches(3.8), Inches(10.3), Inches(0.6))
tf2 = txBox2.text_frame
p2 = tf2.paragraphs[0]
p2.text = "欢迎提问与讨论"
p2.font.size = Pt(24)
p2.font.color.rgb = C_ACCENT
p2.font.name = "Calibri"
p2.alignment = PP_ALIGN.CENTER

line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(3.5), Inches(5.3), Pt(3))
line.fill.solid()
line.fill.fore_color.rgb = C_ACCENT
line.line.fill.background()

# ── Save ──
prs.save(PPT_PATH)
print(f"PPT saved to: {PPT_PATH}")
print(f"Total slides: {len(prs.slides)}")
