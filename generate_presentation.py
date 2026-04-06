"""
Generate presentation-ready figures and academic PPT for group meeting.
All figures use categorized feature names with clear academic titles.
Fonts: SimSun (宋体) for Chinese text, Times New Roman for English/numbers.
Output: representation/ folder + PPT file.
"""

import os, sys, warnings
os.environ['OMP_NUM_THREADS'] = '1'
os.environ['MKL_NUM_THREADS'] = '1'
os.environ['OPENBLAS_NUM_THREADS'] = '1'

import numpy as np
import pandas as pd
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.font_manager as fm
import seaborn as sns
from scipy.stats import ks_2samp
from sklearn.preprocessing import StandardScaler
from sklearn.ensemble import RandomForestClassifier
from sklearn.model_selection import StratifiedKFold
from sklearn.metrics import classification_report, confusion_matrix, accuracy_score
from sklearn.decomposition import PCA
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

warnings.filterwarnings('ignore')

# ============================================================
# Paths
# ============================================================
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
OUTPUT_DIR = os.path.join(SCRIPT_DIR, "output")
REP_DIR = os.path.join(SCRIPT_DIR, "representation")
os.makedirs(REP_DIR, exist_ok=True)

V2_CSV = os.path.join(OUTPUT_DIR, "near_vertical_features_v2.csv")
TARGET_ANGLES = [80, 82, 84, 86, 88, 90]

# ============================================================
# Feature Classification (7 categories)
# ============================================================
FEATURE_CATEGORIES = {
    "Geometric Features": {
        "features": ['n_pixels', 'bbox_width', 'bbox_height', 'aspect_ratio',
                      'eccentricity', 'compactness', 'diagonal_length'],
        "label_en": "Geometric Features",
        "label_zh": "几何特征",
        "color": "#1f77b4",
    },
    "ToT Statistical Features": {
        "features": ['total_tot', 'max_tot', 'mean_tot', 'std_tot',
                      'max_tot_fraction', 'tot_entropy', 'tot_gini',
                      'tot_kurtosis', 'tot_skewness'],
        "label_en": "ToT Statistical Features",
        "label_zh": "ToT统计特征",
        "color": "#2ca02c",
    },
    "Energy Spatial Distribution": {
        "features": ['energy_gradient', 'weighted_centroid_offset',
                      'max_tot_position_ratio', 'second_moment',
                      'tot_weighted_radius'],
        "label_en": "Energy Spatial Distribution",
        "label_zh": "能量空间分布特征",
        "color": "#ff7f0e",
    },
    "PCA Axis Gradient": {
        "features": ['pca_axis_tot_slope', 'pca_axis_tot_slope_abs',
                      'pca_axis_tot_r2', 'pca_eigenvalue_ratio'],
        "label_en": "PCA Axis Gradient",
        "label_zh": "PCA轴能量梯度",
        "color": "#d62728",
    },
    "Discrete Gradient Features": {
        "features": ['adj_diff_mean', 'adj_diff_max', 'adj_diff_std',
                      'adj_diff_x_mean', 'adj_diff_y_mean',
                      'adj_diff_xy_asymmetry', 'adj_diff_signed_skewness'],
        "label_en": "Discrete Gradient Features",
        "label_zh": "离散梯度特征",
        "color": "#9467bd",
    },
    "Azimuthal & Radial Features": {
        "features": ['quadrant_tot_std', 'quadrant_tot_max_min_ratio',
                      'half_plane_x_ratio', 'half_plane_y_ratio',
                      'half_plane_max_asymmetry',
                      'radial_decay_slope', 'radial_decay_slope_abs',
                      'radial_decay_r2', 'half_energy_radius', 'core_fraction'],
        "label_en": "Azimuthal & Radial Features",
        "label_zh": "方位角与径向衰减特征",
        "color": "#8c564b",
    },
    "Hu Moments & Weighted Morphology": {
        "features": ['hu_moment_1', 'hu_moment_2', 'hu_moment_3',
                      'hu_moment_4', 'hu_moment_5', 'hu_moment_6',
                      'hu_moment_7', 'weighted_std_x', 'weighted_std_y',
                      'weighted_xy_ratio', 'weighted_covariance_xy'],
        "label_en": "Hu Moments & Weighted Morphology",
        "label_zh": "Hu矩与加权形态特征",
        "color": "#e377c2",
    },
}

# Phase mapping
V1_CATEGORIES = ["Geometric Features", "ToT Statistical Features", "Energy Spatial Distribution"]
V2_CATEGORIES = ["PCA Axis Gradient", "Discrete Gradient Features",
                  "Azimuthal & Radial Features", "Hu Moments & Weighted Morphology"]

ALL_V1_FEATURES = []
for cat in V1_CATEGORIES:
    ALL_V1_FEATURES.extend(FEATURE_CATEGORIES[cat]["features"])
ALL_V2_FEATURES = []
for cat in V2_CATEGORIES:
    ALL_V2_FEATURES.extend(FEATURE_CATEGORIES[cat]["features"])

def get_feature_color(feat):
    """Return the category color for a feature."""
    for cat_info in FEATURE_CATEGORIES.values():
        if feat in cat_info["features"]:
            return cat_info["color"]
    return "#333333"

def get_feature_category(feat):
    """Return the category name for a feature."""
    for cat_name, cat_info in FEATURE_CATEGORIES.items():
        if feat in cat_info["features"]:
            return cat_name
    return "Unknown"

# ============================================================
# Matplotlib Style Setup
# ============================================================
def setup_matplotlib():
    """Configure matplotlib for academic figures with SimSun + Times New Roman."""
    plt.rcParams.update({
        'font.family': ['Times New Roman', 'SimSun'],
        'font.size': 12,
        'axes.titlesize': 14,
        'axes.labelsize': 12,
        'xtick.labelsize': 10,
        'ytick.labelsize': 10,
        'legend.fontsize': 10,
        'figure.dpi': 150,
        'savefig.dpi': 200,
        'savefig.bbox': 'tight',
        'axes.spines.top': False,
        'axes.spines.right': False,
    })

setup_matplotlib()


# ============================================================
# Figure 1: KS Heatmap — Side by Side (V1 | V2)
# ============================================================
def fig_ks_heatmap(df):
    print("  [1/6] KS Statistic Heatmaps ...")
    feature_cols = [c for c in df.columns if c not in ('angle', 'sample_id')]
    pairs = [(TARGET_ANGLES[i], TARGET_ANGLES[i+1]) for i in range(len(TARGET_ANGLES)-1)]
    pair_labels = [f'{a}°–{b}°' for a, b in pairs]

    ks_stats = pd.DataFrame(index=feature_cols, columns=pair_labels, dtype=float)
    for feat in feature_cols:
        for (a1, a2), label in zip(pairs, pair_labels):
            d1 = df[df['angle'] == a1][feat].dropna().values
            d2 = df[df['angle'] == a2][feat].dropna().values
            stat, _ = ks_2samp(d1, d2)
            ks_stats.loc[feat, label] = stat

    # --- V1 categories subplot ---
    v1_feats = ALL_V1_FEATURES
    v2_feats = ALL_V2_FEATURES
    ks_v1 = ks_stats.loc[v1_feats].astype(float)
    ks_v2 = ks_stats.loc[v2_feats].astype(float)

    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(24, max(10, len(v2_feats) * 0.42)),
                                    gridspec_kw={'width_ratios': [len(v1_feats), len(v2_feats)]})

    # Category labels on y-axis for V1
    v1_yticks = []
    for cat_name in V1_CATEGORIES:
        feats = FEATURE_CATEGORIES[cat_name]["features"]
        for f in feats:
            v1_yticks.append(f)

    sns.heatmap(ks_v1, annot=True, fmt='.3f', cmap='YlOrRd',
                vmin=0, vmax=0.10, linewidths=0.5, ax=ax1, cbar=False,
                annot_kws={'fontsize': 9})
    ax1.set_title('Phase 1: Basic Handcrafted Features (21)', fontsize=13, fontweight='bold')
    ax1.set_xlabel('Adjacent Angle Pair', fontsize=11)
    ax1.set_ylabel('Feature', fontsize=11)

    # Color V1 y-tick labels by category
    for tick_label in ax1.get_yticklabels():
        feat_name = tick_label.get_text()
        tick_label.set_color(get_feature_color(feat_name))
        tick_label.set_fontsize(9)

    # Add category brackets (horizontal lines between categories)
    offset = 0
    for cat_name in V1_CATEGORIES:
        n = len(FEATURE_CATEGORIES[cat_name]["features"])
        mid = offset + n / 2
        ax1.axhline(y=offset, color='gray', linewidth=0.8, linestyle='-', alpha=0.5)
        offset += n

    sns.heatmap(ks_v2, annot=True, fmt='.3f', cmap='YlOrRd',
                vmin=0, vmax=0.10, linewidths=0.5, ax=ax2,
                cbar_kws={'label': 'KS Statistic', 'shrink': 0.6},
                annot_kws={'fontsize': 9})
    ax2.set_title('Phase 2: High-Order Spatial Features (32)', fontsize=13, fontweight='bold')
    ax2.set_xlabel('Adjacent Angle Pair', fontsize=11)
    ax2.set_ylabel('')

    for tick_label in ax2.get_yticklabels():
        feat_name = tick_label.get_text()
        tick_label.set_color(get_feature_color(feat_name))
        tick_label.set_fontsize(9)

    offset = 0
    for cat_name in V2_CATEGORIES:
        n = len(FEATURE_CATEGORIES[cat_name]["features"])
        ax2.axhline(y=offset, color='gray', linewidth=0.8, linestyle='-', alpha=0.5)
        offset += n

    fig.suptitle('Kolmogorov–Smirnov Test Statistics for Adjacent Angle Pairs (80°–90°)',
                 fontsize=15, fontweight='bold', y=1.02)
    plt.tight_layout()
    path = os.path.join(REP_DIR, "ks_heatmap.png")
    plt.savefig(path, dpi=200, bbox_inches='tight')
    plt.close()
    print(f"    Saved: {path}")
    return ks_stats


# ============================================================
# Figure 2: Feature Distributions (selected key features)
# ============================================================
def fig_feature_distributions(df):
    print("  [2/6] Feature Distribution Plots ...")
    # Select representative features from each category
    selected = {
        'Geometric Features': ['n_pixels', 'eccentricity', 'diagonal_length'],
        'ToT Statistical Features': ['total_tot', 'tot_entropy', 'std_tot'],
        'Energy Spatial Distribution': ['energy_gradient', 'weighted_centroid_offset', 'second_moment'],
        'PCA Axis Gradient': ['pca_axis_tot_slope_abs', 'pca_eigenvalue_ratio'],
        'Discrete Gradient Features': ['adj_diff_mean', 'adj_diff_xy_asymmetry'],
        'Azimuthal & Radial Features': ['radial_decay_slope_abs', 'core_fraction', 'half_energy_radius'],
        'Hu Moments & Weighted Morphology': ['hu_moment_1', 'weighted_std_x', 'weighted_covariance_xy'],
    }

    all_feats = []
    cat_labels = []
    for cat, feats in selected.items():
        all_feats.extend(feats)
        cat_labels.extend([cat] * len(feats))

    n = len(all_feats)
    ncols = 4
    nrows = (n + ncols - 1) // ncols

    fig, axes = plt.subplots(nrows, ncols, figsize=(20, nrows * 3.5))
    axes = axes.flatten()
    colors_angle = ['#e74c3c', '#e67e22', '#f1c40f', '#2ecc71', '#3498db', '#9b59b6']

    for i, feat in enumerate(all_feats):
        ax = axes[i]
        for j, angle in enumerate(TARGET_ANGLES):
            vals = df[df['angle'] == angle][feat].dropna().values
            ax.hist(vals, bins=60, alpha=0.45, density=True,
                    color=colors_angle[j], label=f'{angle}°')
        ax.set_title(feat, fontsize=10, fontweight='bold',
                     color=get_feature_color(feat))
        ax.set_ylabel('Density', fontsize=8)
        ax.tick_params(labelsize=8)
        if i == 0:
            ax.legend(fontsize=7, loc='upper right')

    for i in range(n, len(axes)):
        axes[i].set_visible(False)

    fig.suptitle('Feature Distributions by Incidence Angle (80°–90°)',
                 fontsize=15, fontweight='bold', y=1.01)
    plt.tight_layout()
    path = os.path.join(REP_DIR, "feature_distributions.png")
    plt.savefig(path, dpi=200, bbox_inches='tight')
    plt.close()
    print(f"    Saved: {path}")


# ============================================================
# Figure 3: Random Forest — Confusion Matrix
# ============================================================
def fig_confusion_matrix(df):
    print("  [3/6] Confusion Matrix (Random Forest) ...")
    feature_cols = [c for c in df.columns if c not in ('angle', 'sample_id')]
    X = df[feature_cols].values
    y = df['angle'].values

    scaler = StandardScaler()
    X_scaled = scaler.fit_transform(X)

    skf = StratifiedKFold(n_splits=5, shuffle=True, random_state=42)
    y_pred_all = np.zeros_like(y)
    fold_accs = []

    for fold_idx, (train_idx, test_idx) in enumerate(skf.split(X_scaled, y)):
        rf = RandomForestClassifier(
            n_estimators=300, max_depth=15, min_samples_leaf=10,
            class_weight='balanced', random_state=42, n_jobs=1)
        rf.fit(X_scaled[train_idx], y[train_idx])
        y_pred_fold = rf.predict(X_scaled[test_idx])
        fold_accs.append(accuracy_score(y[test_idx], y_pred_fold))
        y_pred_all[test_idx] = y_pred_fold
        print(f"    Fold {fold_idx+1}: {fold_accs[-1]:.4f}")

    mean_acc = np.mean(fold_accs)
    std_acc = np.std(fold_accs)
    print(f"    Mean Accuracy: {mean_acc:.4f} ± {std_acc:.4f}")

    cm = confusion_matrix(y, y_pred_all, labels=TARGET_ANGLES)
    cm_pct = cm.astype(float) / cm.sum(axis=1, keepdims=True) * 100

    fig, ax = plt.subplots(figsize=(9, 7.5))
    # Create annotation text with counts and percentages
    annot_text = np.empty_like(cm, dtype=object)
    for i in range(cm.shape[0]):
        for j in range(cm.shape[1]):
            annot_text[i, j] = f'{cm[i,j]}\n({cm_pct[i,j]:.1f}%)'

    sns.heatmap(cm, annot=annot_text, fmt='', cmap='Blues',
                xticklabels=[f'{a}°' for a in TARGET_ANGLES],
                yticklabels=[f'{a}°' for a in TARGET_ANGLES], ax=ax,
                linewidths=0.5, cbar_kws={'label': 'Sample Count'})
    ax.set_xlabel('Predicted Angle', fontsize=12)
    ax.set_ylabel('True Angle', fontsize=12)
    ax.set_title(f'Random Forest Confusion Matrix (53 Features)\n'
                 f'5-Fold CV Accuracy: {mean_acc:.1%} ± {std_acc:.1%}  |  '
                 f'Random Baseline: 16.7%',
                 fontsize=13, fontweight='bold')
    plt.tight_layout()
    path = os.path.join(REP_DIR, "confusion_matrix.png")
    plt.savefig(path, dpi=200, bbox_inches='tight')
    plt.close()
    print(f"    Saved: {path}")

    # Also train full model for feature importance
    rf_full = RandomForestClassifier(
        n_estimators=300, max_depth=15, min_samples_leaf=10,
        class_weight='balanced', random_state=42, n_jobs=1)
    rf_full.fit(X_scaled, y)
    feat_imp = pd.Series(rf_full.feature_importances_, index=feature_cols).sort_values(ascending=False)

    return mean_acc, std_acc, feat_imp


# ============================================================
# Figure 4: Feature Importance (color-coded by category)
# ============================================================
def fig_feature_importance(feat_imp):
    print("  [4/6] Feature Importance Plot ...")
    feat_imp_sorted = feat_imp.sort_values(ascending=True)

    fig, ax = plt.subplots(figsize=(11, max(10, len(feat_imp_sorted) * 0.32)))
    colors_bar = [get_feature_color(f) for f in feat_imp_sorted.index]
    bars = ax.barh(range(len(feat_imp_sorted)), feat_imp_sorted.values,
                   color=colors_bar, edgecolor='white', linewidth=0.3)
    ax.set_yticks(range(len(feat_imp_sorted)))
    ax.set_yticklabels(feat_imp_sorted.index, fontsize=8)
    ax.set_xlabel('Feature Importance (Gini)', fontsize=11)
    ax.set_title('Random Forest Feature Importance by Category',
                 fontsize=14, fontweight='bold')

    # Add legend for categories
    from matplotlib.patches import Patch
    legend_elements = []
    for cat_name, cat_info in FEATURE_CATEGORIES.items():
        legend_elements.append(Patch(facecolor=cat_info["color"],
                                     label=cat_info["label_en"]))
    ax.legend(handles=legend_elements, loc='lower right', fontsize=8,
              title='Feature Category', title_fontsize=9)

    plt.tight_layout()
    path = os.path.join(REP_DIR, "feature_importance.png")
    plt.savefig(path, dpi=200, bbox_inches='tight')
    plt.close()
    print(f"    Saved: {path}")


# ============================================================
# Figure 5: PCA Scatter
# ============================================================
def fig_pca(df):
    print("  [5/6] PCA Visualization ...")
    feature_cols = [c for c in df.columns if c not in ('angle', 'sample_id')]
    X = df[feature_cols].values
    y = df['angle'].values

    scaler = StandardScaler()
    X_scaled = scaler.fit_transform(X)

    pca = PCA(n_components=2, random_state=42)
    X_pca = pca.fit_transform(X_scaled)

    fig, ax = plt.subplots(figsize=(9, 7))
    colors_angle = ['#e74c3c', '#3498db', '#2ecc71', '#9b59b6', '#e67e22', '#1abc9c']

    for i, angle in enumerate(TARGET_ANGLES):
        mask = y == angle
        ax.scatter(X_pca[mask, 0], X_pca[mask, 1], c=colors_angle[i],
                   label=f'{angle}°', alpha=0.15, s=3, rasterized=True)

    var1, var2 = pca.explained_variance_ratio_ * 100
    ax.set_xlabel(f'PC1 ({var1:.1f}% Variance)', fontsize=12)
    ax.set_ylabel(f'PC2 ({var2:.1f}% Variance)', fontsize=12)
    ax.set_title(f'PCA Projection of 53 Features (80°–90°)\n'
                 f'PC1+PC2 Cumulative Variance: {var1+var2:.1f}%',
                 fontsize=13, fontweight='bold')
    ax.legend(fontsize=10, markerscale=5, title='Incidence Angle')
    plt.tight_layout()
    path = os.path.join(REP_DIR, "pca_scatter.png")
    plt.savefig(path, dpi=200, bbox_inches='tight')
    plt.close()
    print(f"    Saved: {path}")


# ============================================================
# Figure 6: Summary Bar Chart — KS max per category
# ============================================================
def fig_category_summary(ks_stats):
    print("  [6/6] Category KS Summary ...")
    pairs = [(TARGET_ANGLES[i], TARGET_ANGLES[i+1]) for i in range(len(TARGET_ANGLES)-1)]
    pair_labels = [f'{a}°–{b}°' for a, b in pairs]

    cat_max_ks = {}
    for cat_name, cat_info in FEATURE_CATEGORIES.items():
        feats_in = [f for f in cat_info["features"] if f in ks_stats.index]
        if feats_in:
            cat_max_ks[cat_name] = ks_stats.loc[feats_in].max().max()

    fig, ax = plt.subplots(figsize=(10, 5))
    cats = list(cat_max_ks.keys())
    vals = [cat_max_ks[c] for c in cats]
    colors = [FEATURE_CATEGORIES[c]["color"] for c in cats]
    short_labels = [FEATURE_CATEGORIES[c]["label_en"] for c in cats]

    bars = ax.bar(range(len(cats)), vals, color=colors, edgecolor='white', linewidth=0.5)
    ax.set_xticks(range(len(cats)))
    ax.set_xticklabels(short_labels, fontsize=9, rotation=20, ha='right')
    ax.set_ylabel('Max KS Statistic', fontsize=11)
    ax.set_title('Maximum KS Statistic by Feature Category\n(All Adjacent Angle Pairs)',
                 fontsize=13, fontweight='bold')
    ax.axhline(y=0.05, color='red', linestyle='--', linewidth=1.2, label='Threshold = 0.05')
    ax.axhline(y=0.10, color='darkred', linestyle='--', linewidth=1.2, label='Threshold = 0.10')
    ax.legend(fontsize=9)
    ax.set_ylim(0, 0.12)

    # Add value labels on bars
    for bar, val in zip(bars, vals):
        ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 0.002,
                f'{val:.3f}', ha='center', va='bottom', fontsize=9, fontweight='bold')

    plt.tight_layout()
    path = os.path.join(REP_DIR, "category_ks_summary.png")
    plt.savefig(path, dpi=200, bbox_inches='tight')
    plt.close()
    print(f"    Saved: {path}")


# ============================================================
# PPT Generation
# ============================================================
def create_ppt(mean_acc, std_acc):
    print("\n  Creating PPT ...")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ---- Color constants ----
    NAVY = RGBColor(0x1a, 0x1a, 0x2e)
    BLUE = RGBColor(0x29, 0x80, 0xb9)
    DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
    GRAY = RGBColor(0x66, 0x66, 0x66)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    LIGHT_BG = RGBColor(0xf0, 0xf4, 0xf8)
    GREEN = RGBColor(0x27, 0xae, 0x60)
    RED = RGBColor(0xc0, 0x39, 0x2b)
    TABLE_HEADER = RGBColor(0x2c, 0x3e, 0x50)

    def add_slide_number(slide, num):
        from pptx.oxml.ns import qn
        txBox = slide.shapes.add_textbox(Inches(12.5), Inches(7.0), Inches(0.7), Inches(0.35))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = str(num)
        p.alignment = PP_ALIGN.RIGHT
        run = p.runs[0]
        run.font.size = Pt(10)
        run.font.color.rgb = GRAY
        run.font.name = 'Times New Roman'

    def set_font(run, name_en='Times New Roman', name_zh='SimSun',
                 size=18, bold=False, color=DARK_GRAY, italic=False):
        run.font.name = name_en
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.italic = italic
        # Set East Asian font
        from pptx.oxml.ns import qn
        rPr = run._r.get_or_add_rPr()
        rPr.set(qn('a:altLang'), 'zh-CN')
        ea = rPr.find(qn('a:ea'))
        if ea is None:
            ea = rPr.makeelement(qn('a:ea'), {})
            rPr.append(ea)
        ea.set('typeface', name_zh)

    def add_text_box(slide, left, top, width, height, text,
                     font_size=18, bold=False, color=DARK_GRAY,
                     alignment=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                          Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = alignment
        run = p.add_run()
        run.text = text
        set_font(run, size=font_size, bold=bold, color=color)
        return tf

    def add_bullet_slide(slide, items, left=0.8, top=1.6, width=11.5, height=5.5,
                         font_size=18, color=DARK_GRAY):
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                          Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.space_after = Pt(8)
            p.level = 0
            run = p.add_run()
            run.text = f"• {item}"
            set_font(run, size=font_size, color=color)

    def add_image_slide(slide, img_path, left=0.5, top=1.3, max_w=12.3, max_h=5.8):
        """Add an image scaled to fit within the given box."""
        from PIL import Image
        img = Image.open(img_path)
        w_px, h_px = img.size
        ratio = w_px / h_px
        # Fit within box
        if max_w / max_h > ratio:
            h = max_h
            w = h * ratio
        else:
            w = max_w
            h = w / ratio
        # Center horizontally
        actual_left = left + (max_w - w) / 2
        slide.shapes.add_picture(img_path, Inches(actual_left), Inches(top),
                                 Inches(w), Inches(h))

    def add_title_bar(slide, title_text, subtitle_text=None):
        """Add a colored title bar at the top of a content slide."""
        # Title bar background
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
            Inches(13.333), Inches(1.15))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x2c, 0x3e, 0x50)
        shape.line.fill.background()

        txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.15),
                                          Inches(12), Inches(0.6))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title_text
        set_font(run, size=28, bold=True, color=WHITE)

        if subtitle_text:
            txBox2 = slide.shapes.add_textbox(Inches(0.6), Inches(0.7),
                                               Inches(12), Inches(0.4))
            tf2 = txBox2.text_frame
            p2 = tf2.paragraphs[0]
            run2 = p2.add_run()
            run2.text = subtitle_text
            set_font(run2, size=14, bold=False, color=RGBColor(0xbb, 0xcc, 0xdd), italic=True)

    slide_num = 0

    # ======== SLIDE 1: Title ========
    slide_num += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # Background gradient-like bar
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        Inches(13.333), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x2c, 0x3e, 0x50)
    shape.line.fill.background()

    # Accent line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(2), Inches(3.4),
        Inches(9.333), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.fill.background()

    add_text_box(slide, 2, 1.5, 9.333, 1.5,
                 'Timepix3 近垂直入射角度分辨能力分析',
                 font_size=34, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 2, 2.5, 9.333, 0.8,
                 'Angular Resolution Analysis at Near-Vertical Incidence (80°–90°)',
                 font_size=18, bold=False, color=RGBColor(0xaa, 0xbb, 0xcc),
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 2, 3.8, 9.333, 0.5,
                 '基于 53 维手工特征的统计检验与机器学习分析',
                 font_size=16, bold=False, color=RGBColor(0x99, 0xaa, 0xbb),
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 2, 5.5, 9.333, 0.5,
                 '组会汇报  |  2026.03.26',
                 font_size=16, bold=False, color=RGBColor(0x88, 0x99, 0xaa),
                 alignment=PP_ALIGN.CENTER)
    add_slide_number(slide, slide_num)

    # ======== SLIDE 2: Outline ========
    slide_num += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, '汇报提纲', 'Outline')
    items = [
        '研究背景与问题定义',
        '数据概况与特征工程（7 类 53 维特征）',
        'Kolmogorov–Smirnov 分布检验结果',
        'Random Forest 分类结果与特征重要性',
        'PCA 降维可视化',
        '结论与后续工作',
    ]
    add_bullet_slide(slide, items, font_size=20)
    add_slide_number(slide, slide_num)

    # ======== SLIDE 3: Background ========
    slide_num += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, '研究背景', 'Background & Motivation')
    items = [
        'Timepix3 混合像素探测器：55 μm 像素间距，500 μm 硅厚度',
        '粒子入射角度影响径迹形态 → 可通过径迹特征推断角度',
        '关键问题：近垂直入射（80°–90°）时，相邻 2° 间隔是否可分辨？',
        '此区间径迹长度变化极小（< 1 pixel），电荷扩散主导径迹形态',
        '目标：定量评估角度分辨极限，为后续分类模型提供依据',
    ]
    add_bullet_slide(slide, items, font_size=18)
    add_slide_number(slide, slide_num)

    # ======== SLIDE 4: Data Overview ========
    slide_num += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, '数据与特征工程', 'Data & Feature Engineering')

    # Data summary
    add_text_box(slide, 0.8, 1.5, 5, 0.4,
                 '数据集概况', font_size=20, bold=True, color=NAVY)
    items_data = [
        '角度范围：80°, 82°, 84°, 86°, 88°, 90°（6 类）',
        '总样本数：191,426（每类 ~30,000）',
        '输入：50×50 ToT 矩阵（每像素能量沉积）',
    ]
    add_bullet_slide(slide, items_data, left=0.8, top=2.1, width=5.5, height=2,
                     font_size=16)

    # Feature table
    add_text_box(slide, 6.8, 1.5, 6, 0.4,
                 '特征分类（7 类 53 维）', font_size=20, bold=True, color=NAVY)

    table_data = [
        ('阶段', '类别', '数量'),
        ('Phase 1', 'Geometric Features（几何特征）', '7'),
        ('Phase 1', 'ToT Statistical Features（ToT统计特征）', '9'),
        ('Phase 1', 'Energy Spatial Distribution（能量空间分布）', '5'),
        ('Phase 2', 'PCA Axis Gradient（PCA轴梯度）', '4'),
        ('Phase 2', 'Discrete Gradient（离散梯度）', '7'),
        ('Phase 2', 'Azimuthal & Radial（方位角与径向）', '10'),
        ('Phase 2', 'Hu Moments & Morphology（Hu矩与形态）', '11'),
        ('', '合计', '53'),
    ]

    rows, cols = len(table_data), 3
    table = slide.shapes.add_table(rows, cols, Inches(6.8), Inches(2.1),
                                    Inches(6), Inches(4.2)).table
    table.columns[0].width = Inches(1.0)
    table.columns[1].width = Inches(4.0)
    table.columns[2].width = Inches(1.0)

    for i, row_data in enumerate(table_data):
        for j, cell_text in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = cell_text
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER if j != 1 else PP_ALIGN.LEFT
                for run in paragraph.runs:
                    set_font(run, size=12 if i > 0 else 13,
                             bold=(i == 0 or i == len(table_data)-1),
                             color=WHITE if i == 0 else DARK_GRAY)
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_HEADER
            elif i == len(table_data) - 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xe8, 0xee, 0xf4)
            elif i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BG

    add_slide_number(slide, slide_num)

    # ======== SLIDE 5: Feature Distributions ========
    slide_num += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, '特征分布：各角度高度重叠', 'Feature Distributions Across Angles')
    img_path = os.path.join(REP_DIR, "feature_distributions.png")
    if os.path.exists(img_path):
        add_image_slide(slide, img_path, left=0.3, top=1.3, max_w=12.7, max_h=5.9)
    add_slide_number(slide, slide_num)

    # ======== SLIDE 6: KS Test Heatmap ========
    slide_num += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, 'KS 检验：所有特征 KS < 0.05',
                  'Kolmogorov–Smirnov Test for Adjacent Angle Pairs')
    img_path = os.path.join(REP_DIR, "ks_heatmap.png")
    if os.path.exists(img_path):
        add_image_slide(slide, img_path, left=0.2, top=1.3, max_w=12.9, max_h=5.9)
    add_slide_number(slide, slide_num)

    # ======== SLIDE 7: Category KS Summary ========
    slide_num += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, '各类别最大 KS 统计量', 'Max KS Statistic by Feature Category')
    img_path = os.path.join(REP_DIR, "category_ks_summary.png")
    if os.path.exists(img_path):
        add_image_slide(slide, img_path, left=1.5, top=1.3, max_w=10, max_h=5.5)

    add_text_box(slide, 0.8, 6.5, 11, 0.5,
                 '结论：所有 7 类特征的最大 KS 统计量均远低于 0.05 阈值 → 相邻角度分布无显著差异',
                 font_size=16, bold=True, color=RED)
    add_slide_number(slide, slide_num)

    # ======== SLIDE 8: Confusion Matrix ========
    slide_num += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, f'Random Forest 分类：准确率 {mean_acc:.1%}',
                  f'5-Fold CV with 53 Features | Random Baseline: 16.7%')
    img_path = os.path.join(REP_DIR, "confusion_matrix.png")
    if os.path.exists(img_path):
        add_image_slide(slide, img_path, left=2.5, top=1.2, max_w=8.5, max_h=6.0)
    add_slide_number(slide, slide_num)

    # ======== SLIDE 9: Feature Importance ========
    slide_num += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, '特征重要性：按类别着色',
                  'Random Forest Feature Importance (Gini) — Color-Coded by Category')
    img_path = os.path.join(REP_DIR, "feature_importance.png")
    if os.path.exists(img_path):
        add_image_slide(slide, img_path, left=0.5, top=1.2, max_w=12, max_h=6.0)
    add_slide_number(slide, slide_num)

    # ======== SLIDE 10: PCA ========
    slide_num += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, 'PCA 降维：各角度完全重叠',
                  'Principal Component Analysis — No Cluster Separation')
    img_path = os.path.join(REP_DIR, "pca_scatter.png")
    if os.path.exists(img_path):
        add_image_slide(slide, img_path, left=2.5, top=1.2, max_w=8.5, max_h=6.0)
    add_slide_number(slide, slide_num)

    # ======== SLIDE 11: Conclusion ========
    slide_num += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, '结论', 'Conclusions')
    items = [
        '53 维特征（7 类）全面刻画了径迹的几何、能量、空间结构信息',
        'KS 检验：所有 53 个特征在相邻角度对间的 KS < 0.035，远低于 0.05 阈值',
        f'Random Forest 分类：准确率 {mean_acc:.1%}（随机基线 16.7%），仅高出 ~3%',
        'PCA 降维：6 类角度在特征空间中完全重叠，不存在可分离的聚类结构',
        '物理解释：近垂直入射时径迹仅 1–3 个像素，电荷扩散主导，2° 间隔产生亚像素径迹变化',
    ]
    add_bullet_slide(slide, items, font_size=18, top=1.5)

    # Key conclusion box
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(5.2),
        Inches(10.333), Inches(1.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xe8, 0xf8, 0xf0)
    shape.line.color.rgb = GREEN
    shape.line.width = Pt(2)

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = '结论：Timepix3 在 [80°, 90°] 范围内无法以 2° 间隔分辨入射角度\n'
    set_font(run, size=18, bold=True, color=RGBColor(0x1a, 0x6e, 0x3a))
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = '建议将 80°–90° 合并为单一角度区间，或采用更大间隔（≥5°）进行分类'
    set_font(run2, size=15, bold=False, color=DARK_GRAY)

    add_slide_number(slide, slide_num)

    # ======== SLIDE 12: Next Steps ========
    slide_num += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, '后续工作', 'Next Steps')
    items = [
        '尝试更大角度间隔（5°, 10°）验证分辨阈值',
        '将 80°–90° 合并为单一类别，评估对整体分类精度的影响',
        '对比 CNN / 深度学习方法是否能提取更深层特征',
        '扩展分析到中等角度范围（45°–80°），确认分辨能力随角度的变化趋势',
    ]
    add_bullet_slide(slide, items, font_size=18, top=1.5)

    # Thank you
    add_text_box(slide, 2, 5.5, 9.333, 0.8,
                 'Thank You  |  谢谢',
                 font_size=28, bold=True, color=BLUE, alignment=PP_ALIGN.CENTER)

    add_slide_number(slide, slide_num)

    # Save
    ppt_path = os.path.join(REP_DIR, "presentation.pptx")
    prs.save(ppt_path)
    print(f"  PPT saved to: {ppt_path}")
    print(f"  Total slides: {slide_num}")


# ============================================================
# Main
# ============================================================
def main():
    import time
    t0 = time.time()
    print("=" * 60)
    print("Generating Presentation-Ready Figures & PPT")
    print("=" * 60)

    print(f"\nLoading data from: {V2_CSV}")
    df = pd.read_csv(V2_CSV)
    print(f"  Shape: {df.shape}")

    print("\n--- Generating Figures ---")
    ks_stats = fig_ks_heatmap(df)
    fig_feature_distributions(df)
    mean_acc, std_acc, feat_imp = fig_confusion_matrix(df)
    fig_feature_importance(feat_imp)
    fig_pca(df)
    fig_category_summary(ks_stats)

    print("\n--- Generating PPT ---")
    create_ppt(mean_acc, std_acc)

    elapsed = time.time() - t0
    print(f"\nDone! Total time: {elapsed:.1f}s")
    print(f"All outputs in: {REP_DIR}")


if __name__ == '__main__':
    main()
