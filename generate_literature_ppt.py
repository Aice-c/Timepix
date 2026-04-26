"""
Generate academic presentation for group meeting:
Literature Review — Tetris-inspired Detector (Okabe et al., Nature Comm. 2024)
and Inspirations for Timepix3 Charged Particle Angle Identification
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# ── Color palette ──
C_TITLE   = RGBColor(0x2c, 0x3e, 0x50)
C_ACCENT  = RGBColor(0x29, 0x80, 0xb9)
C_BODY    = RGBColor(0x33, 0x33, 0x33)
C_CAPTION = RGBColor(0x66, 0x66, 0x66)
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_BG_ALT  = RGBColor(0xF0, 0xF4, 0xF8)
C_GREEN   = RGBColor(0x27, 0xAE, 0x60)
C_RED     = RGBColor(0xC0, 0x39, 0x2B)
C_ORANGE  = RGBColor(0xED, 0x7D, 0x31)
C_LIGHT_LINE = RGBColor(0xBD, 0xC3, 0xC7)

OUTPUT_DIR = r"d:\Project\Timepix\output"
os.makedirs(OUTPUT_DIR, exist_ok=True)
PPT_PATH = os.path.join(OUTPUT_DIR, "literature_review_presentation.pptx")

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)


# ══════════════════════════════════════════════════════════════
#  Helper functions (reused from existing codebase)
# ══════════════════════════════════════════════════════════════

def add_slide():
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def add_title_bar(slide, text, y=Inches(0.3), font_size=Pt(32)):
    txBox = slide.shapes.add_textbox(Inches(0.8), y, Inches(11.7), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.bold = True
    p.font.color.rgb = C_TITLE
    p.font.name = "Calibri"
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), y + Inches(0.85),
        Inches(11.7), Pt(3)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = C_ACCENT
    line.line.fill.background()
    return tf


def add_bullet_slide(slide, title, bullets, start_top=Inches(1.5)):
    add_title_bar(slide, title)
    txBox = slide.shapes.add_textbox(Inches(0.8), start_top, Inches(11.7), Inches(5.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, level) in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.level = level
        p.font.size = Pt(22) if level == 0 else Pt(19)
        p.font.color.rgb = C_BODY if level == 0 else C_CAPTION
        p.font.name = "Calibri"
        p.space_before = Pt(10) if level == 0 else Pt(4)
        p.space_after = Pt(4)
    return tf


def add_body_text(slide, text, left=Inches(0.8), top=Inches(1.5),
                  width=Inches(11.7), height=Inches(5.2), font_size=Pt(20)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_text in enumerate(text.strip().split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line_text
        p.font.size = font_size
        p.font.color.rgb = C_BODY
        p.font.name = "Calibri"
        p.space_after = Pt(6)
    return tf


def add_table(slide, data, left=Inches(0.8), top=Inches(1.5),
              col_widths=None, font_size=Pt(16)):
    rows, cols = len(data), len(data[0])
    if col_widths is None:
        w = Inches(11.7)
        col_widths = [w // cols] * cols
    table_shape = slide.shapes.add_table(rows, cols, left, top,
                                          sum(col_widths), Inches(0.45 * rows))
    table = table_shape.table
    for ci, cw in enumerate(col_widths):
        table.columns[ci].width = cw
    for ri, row_data in enumerate(data):
        for ci, cell_text in enumerate(row_data):
            cell = table.cell(ri, ci)
            cell.text = str(cell_text)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = font_size
                paragraph.font.name = "Calibri"
                paragraph.alignment = PP_ALIGN.CENTER
                if ri == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = C_WHITE
                else:
                    paragraph.font.color.rgb = C_BODY
            if ri == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_TITLE
            elif ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_BG_ALT
    return table


def add_slide_number(slide, num):
    txBox = slide.shapes.add_textbox(Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = str(num)
    p.font.size = Pt(12)
    p.font.color.rgb = C_CAPTION
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.RIGHT


def add_two_column_text(slide, title, left_title, left_bullets, right_title, right_bullets):
    """Add a slide with two text columns."""
    add_title_bar(slide, title)
    col_w = Inches(5.5)
    # Left column title
    txL_title = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), col_w, Inches(0.5))
    tf = txL_title.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = C_ACCENT
    p.font.name = "Calibri"
    # Left column body
    txL = slide.shapes.add_textbox(Inches(0.8), Inches(2.1), col_w, Inches(4.8))
    tf = txL.text_frame
    tf.word_wrap = True
    for i, (text, level) in enumerate(left_bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.level = level
        p.font.size = Pt(19) if level == 0 else Pt(16)
        p.font.color.rgb = C_BODY
        p.font.name = "Calibri"
        p.space_before = Pt(6)
        p.space_after = Pt(3)
    # Right column title
    txR_title = slide.shapes.add_textbox(Inches(6.8), Inches(1.5), col_w, Inches(0.5))
    tf = txR_title.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = C_ORANGE
    p.font.name = "Calibri"
    # Right column body
    txR = slide.shapes.add_textbox(Inches(6.8), Inches(2.1), col_w, Inches(4.8))
    tf = txR.text_frame
    tf.word_wrap = True
    for i, (text, level) in enumerate(right_bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.level = level
        p.font.size = Pt(19) if level == 0 else Pt(16)
        p.font.color.rgb = C_BODY
        p.font.name = "Calibri"
        p.space_before = Pt(6)
        p.space_after = Pt(3)
    # Vertical divider
    div = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(6.5), Inches(1.5), Pt(2), Inches(5.0)
    )
    div.fill.solid()
    div.fill.fore_color.rgb = C_LIGHT_LINE
    div.line.fill.background()


def add_problem_solution_slide(slide, title, problem_text, solution_text, reason_text):
    """Slide with Problem → Solution → Why it works structure."""
    add_title_bar(slide, title)
    labels = ["问题", "启发方案", "为什么有效"]
    colors = [C_RED, C_GREEN, C_ACCENT]
    texts = [problem_text, solution_text, reason_text]
    top = Inches(1.5)
    for label, color, text in zip(labels, colors, texts):
        # Label box
        lbl = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top, Inches(1.8), Inches(0.5)
        )
        lbl.fill.solid()
        lbl.fill.fore_color.rgb = color
        lbl.line.fill.background()
        for paragraph in lbl.text_frame.paragraphs:
            paragraph.text = label
            paragraph.font.size = Pt(16)
            paragraph.font.bold = True
            paragraph.font.color.rgb = C_WHITE
            paragraph.font.name = "Calibri"
            paragraph.alignment = PP_ALIGN.CENTER
        lbl.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        # Content text
        txBox = slide.shapes.add_textbox(Inches(3.0), top, Inches(9.5), Inches(1.4))
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, line_text in enumerate(text.strip().split("\n")):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line_text
            p.font.size = Pt(18)
            p.font.color.rgb = C_BODY
            p.font.name = "Calibri"
            p.space_after = Pt(3)
        top += Inches(1.7)


# ══════════════════════════════════════════════════════════════
#  SLIDE 1: Title Slide
# ══════════════════════════════════════════════════════════════
slide_num = 1
s = add_slide()
txBox = s.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(10.9), Inches(1.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "文献解读与方法启发"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = C_TITLE
p.font.name = "Calibri"
p.alignment = PP_ALIGN.CENTER

txBox2 = s.shapes.add_textbox(Inches(1.2), Inches(3.4), Inches(10.9), Inches(1.2))
tf2 = txBox2.text_frame
tf2.word_wrap = True
p2 = tf2.paragraphs[0]
p2.text = "Tetris-inspired Detector with Neural Network for Radiation Mapping"
p2.font.size = Pt(22)
p2.font.color.rgb = C_ACCENT
p2.font.name = "Calibri"
p2.alignment = PP_ALIGN.CENTER
p3 = tf2.add_paragraph()
p3.text = "Okabe et al., Nature Communications, 2024"
p3.font.size = Pt(18)
p3.font.color.rgb = C_CAPTION
p3.font.name = "Calibri"
p3.alignment = PP_ALIGN.CENTER

txBox3 = s.shapes.add_textbox(Inches(1.2), Inches(5.2), Inches(10.9), Inches(0.5))
tf3 = txBox3.text_frame
p4 = tf3.paragraphs[0]
p4.text = "组会汇报  |  2026年4月"
p4.font.size = Pt(18)
p4.font.color.rgb = C_CAPTION
p4.font.name = "Calibri"
p4.alignment = PP_ALIGN.CENTER


# ══════════════════════════════════════════════════════════════
#  SLIDE 2: Outline
# ══════════════════════════════════════════════════════════════
slide_num += 1
s = add_slide()
add_bullet_slide(s, "汇报提纲", [
    ("一、论文概述 — 做了什么？", 0),
    ("二、论文关键技术 — 怎么做的？", 0),
    ("三、与本课题的对比 — 物理场景差异", 0),
    ("四、与本课题的对比 — 技术方案差异", 0),
    ("五、启发一：Wasserstein 损失函数", 0),
    ("六、启发二：高斯软标签", 0),
    ("七、启发三：浅层网络架构适配", 0),
    ("八、前期分析结论：80°–90° 物理不可分辨", 0),
    ("九、下一步实验计划", 0),
    ("十、总结", 0),
])
add_slide_number(s, slide_num)


# ══════════════════════════════════════════════════════════════
#  SLIDE 3: Paper Overview
# ══════════════════════════════════════════════════════════════
slide_num += 1
s = add_slide()
add_bullet_slide(s, "一、论文概述", [
    ("研究目标：用最少的探测器像素实现高精度 γ 射线方向探测", 0),
    ("仅用 4 个 CZT 晶体像素（1cm×1cm），像素间填充 1mm 铅屏蔽材料", 1),
    ("核心创新：Tetris 形状探测器 + 神经网络分析", 0),
    ("探索了 2×2 方形、S形、J形、T形等 Tetromino 构型", 1),
    ("S形探测器表现最佳，平均方向预测误差 ~0.95°", 1),
    ("数据来源：OpenMC 蒙特卡罗模拟", 0),
    ("0.5 MeV γ 射线，源距 20–500 cm，全 360° 方位角扫描", 1),
    ("应用扩展：移动探测器 + MAP 估计实现辐射源空间定位", 0),
    ("在真实环境中用 Cs-137 源验证了定位能力", 1),
])
add_slide_number(s, slide_num)


# ══════════════════════════════════════════════════════════════
#  SLIDE 4: Paper Key Techniques
# ══════════════════════════════════════════════════════════════
slide_num += 1
s = add_slide()
add_title_bar(s, "二、论文关键技术")

# Three key technique boxes
techniques = [
    ("① 像素间铅填充增强对比度",
     "不同方向的 γ 射线被铅屏蔽遮挡的程度不同\n"
     "→ 4个像素的计数比产生方向依赖信号\n"
     "→ 人为制造方向敏感性"),
    ("② Filter Layer + 1D U-Net 架构",
     "Filter Layer：用高质量MC模拟的响应函数初始化\n"
     "→ 物理先验注入网络（类似模板匹配）\n"
     "U-Net：像素级精细回归角度概率分布"),
    ("③ Wasserstein 距离损失函数",
     "输出 64 维角度概率分布（非离散类别）\n"
     "用 Earth Mover's Distance 度量分布差异\n"
     "训练用 2 阶 W₂（收敛快），评估用 1 阶 W₁（=角度误差）"),
]

top = Inches(1.5)
for title_text, body_text in techniques:
    # Title
    txT = s.shapes.add_textbox(Inches(0.8), top, Inches(11.7), Inches(0.45))
    tf = txT.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = C_ACCENT
    p.font.name = "Calibri"
    # Body
    txB = s.shapes.add_textbox(Inches(1.2), top + Inches(0.45), Inches(11.3), Inches(1.3))
    tf = txB.text_frame
    tf.word_wrap = True
    for i, line_text in enumerate(body_text.split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line_text
        p.font.size = Pt(17)
        p.font.color.rgb = C_BODY
        p.font.name = "Calibri"
        p.space_after = Pt(2)
    top += Inches(1.8)

add_slide_number(s, slide_num)


# ══════════════════════════════════════════════════════════════
#  SLIDE 5: Comparison — Physics Scenario
# ══════════════════════════════════════════════════════════════
slide_num += 1
s = add_slide()
add_title_bar(s, "三、与本课题的对比 — 物理场景")
add_table(s, [
    ["对比维度", "Okabe et al. (2024)", "本课题"],
    ["探测对象", "γ 射线 (0.5 MeV)", "带电粒子: α (Am-241) / 质子"],
    ["方向信息来源", "宏观几何遮挡\n(像素间铅屏蔽)", "微观径迹形态\n(粒子在硅中的能量沉积)"],
    ["角度定义", "方位角 θ ∈ [0°, 360°)\n辐射源围绕探测器的方向", "极角 θ ∈ [0°, 90°]\n粒子穿入探测器的倾斜角"],
    ["探测器", "4个 CZT 晶体 (1cm×1cm)\nTetris 形状排列", "Timepix3 (256×256 像素)\n55μm 间距，单层硅传感器"],
    ["信息承载方式", "4个像素的计数值之比\n(宏观统计量)", "~12个活跃像素的 ToT 分布\n(微观图像)"],
], left=Inches(0.5), top=Inches(1.5),
   col_widths=[Inches(2.2), Inches(4.8), Inches(4.8)],
   font_size=Pt(14))
add_slide_number(s, slide_num)


# ══════════════════════════════════════════════════════════════
#  SLIDE 6: Comparison — Technical Approach
# ══════════════════════════════════════════════════════════════
slide_num += 1
s = add_slide()
add_title_bar(s, "四、与本课题的对比 — 技术方案")
add_table(s, [
    ["对比维度", "Okabe et al. (2024)", "本课题（改进前）"],
    ["输入数据", "2×2 或 2×3 矩阵 (4~6个值)\n信息密度 100%", "50×50 矩阵 (2500个值)\n仅 ~12 个非零，密度 0.48%"],
    ["数据来源", "Monte Carlo 模拟\n可无限生成，无噪声", "真实实验数据\n样本有限，含统计涨落"],
    ["任务定义", "角度分布回归\n输出 64 维概率分布", "角度分类\n输出离散类别标签"],
    ["网络架构", "定制 Filter Layer + 1D U-Net\n为 4 像素信号量身设计", "ResNet-18 (ImageNet预设)\n5次下采样，不适合稀疏输入"],
    ["损失函数", "Wasserstein 距离\n感知角度有序性", "交叉熵\n所有错误类别等价惩罚"],
],  left=Inches(0.5), top=Inches(1.5),
    col_widths=[Inches(2.2), Inches(4.8), Inches(4.8)],
    font_size=Pt(14))
add_slide_number(s, slide_num)


# ══════════════════════════════════════════════════════════════
#  SLIDE 7: Inspiration 1 — Wasserstein Loss
# ══════════════════════════════════════════════════════════════
slide_num += 1
s = add_slide()
add_problem_solution_slide(s,
    "五、启发一：Wasserstein 损失函数",
    problem_text=(
        "交叉熵损失不感知角度有序性：\n"
        "真实=80°时，预测82°（差2°）和预测90°（差10°）的损失相同"
    ),
    solution_text=(
        "引入一维 Wasserstein 距离（Earth Mover's Distance）：\n"
        "W₁ = Σ|CDF(预测) - CDF(目标)|，通过 CDF 差值自动编码有序性"
    ),
    reason_text=(
        "预测偏到相邻角度 → CDF 差值小 → 损失小\n"
        "预测偏到远处角度 → CDF 差值大 → 损失大\n"
        "模型学到「角度是有顺序的」，梯度信号与物理意义一致"
    ),
)
add_slide_number(s, slide_num)


# ══════════════════════════════════════════════════════════════
#  SLIDE 8: Inspiration 2 — Gaussian Soft Labels
# ══════════════════════════════════════════════════════════════
slide_num += 1
s = add_slide()
add_problem_solution_slide(s,
    "六、启发二：高斯软标签",
    problem_text=(
        "One-hot 标签 + 相邻角度径迹高度相似 → 梯度矛盾：\n"
        "84°样本要求 p(84°)=1，但82°样本径迹几乎一样却要求 p(82°)=1"
    ),
    solution_text=(
        "将 one-hot 标签替换为以真实角度为中心的高斯分布：\n"
        "y_k = (1/Z) · exp(-(θ_k - θ_true)² / 2σ²), σ 反映物理分辨力"
    ),
    reason_text=(
        "相似样本 → 相似的标签分布 → 梯度不再矛盾\n"
        "配合 Wasserstein 损失：目标分布是平滑曲线，训练更加稳定\n"
        "σ 可调：低角度用小 σ（精确），高角度用大 σ（模糊）"
    ),
)
add_slide_number(s, slide_num)


# ══════════════════════════════════════════════════════════════
#  SLIDE 9: Inspiration 3 — Shallow Network
# ══════════════════════════════════════════════════════════════
slide_num += 1
s = add_slide()
add_title_bar(s, "七、启发三：浅层网络架构适配")

add_two_column_text(s, "",
    left_title="ResNet-18 的问题",
    left_bullets=[
        ("为 224×224 密集图像设计，5次下采样", 0),
        ("50×50 输入，活跃区域仅 ~7×7", 0),
        ("下采样链中信息快速丢失：", 0),
        ("  50×50 → 25×25 → 13×13 → 7×7 → 4×4 → 2×2", 1),
        ("  活跃区域：7×7 → 4×4 → 2×2 → 1×1 → 消失", 1),
        ("从 layer2 开始空间结构完全丢失", 0),
        ("网络后半段在处理空信号", 0),
    ],
    right_title="ShallowResNet 改进",
    right_bullets=[
        ("论文启发：用极简架构处理极简数据", 0),
        ("论文仅 4 个输入值 → 定制 Filter+U-Net", 1),
        ("仅 1 次下采样，保留空间分辨率：", 0),
        ("  50×50 → 50×50 → 25×25 → 25×25 → GAP", 1),
        ("  活跃区域始终 ≥ 4×4（有空间结构）", 1),
        ("3-4 层卷积 → 感受野 7×7-9×9", 0),
        ("刚好覆盖团簇，不过度设计", 1),
        ("参数量：~1M vs ResNet-18 ~11M", 0),
    ],
)
add_slide_number(s, slide_num)


# ══════════════════════════════════════════════════════════════
#  SLIDE 10: 80-90 Indistinguishability
# ══════════════════════════════════════════════════════════════
slide_num += 1
s = add_slide()
add_title_bar(s, "八、前期分析结论：80°–90° 物理不可分辨")

# Left: analysis results
add_body_text(s, (
    "已完成两阶段数据分析（基础特征 + 高阶空间特征）：\n"
    "\n"
    "• 所有相邻角度对的 KS 统计量 < 0.035（阈值 0.10）\n"
    "• 随机森林准确率 19.0%（6类随机基线 16.7%）\n"
    "• 高阶空间特征（梯度/方位角/Hu矩）仍无效 → RF ≈ 22%\n"
    "• Wasserstein + ShallowResNet → 15%（仍为随机水平）\n"
    "\n"
    "物理原因：\n"
    "• 80°–90° 投影长度差异 < 1 个像素（55μm）\n"
    "• 电荷扩散尺度（~3-4像素）>> 角度差异尺度\n"
    "• 所有可观测量变化仅 4%–5%，淹没在统计涨落中\n"
    "\n"
    "结论：2° 间隔不可区分是探测器固有物理极限\n"
    "→ 非模型能力问题，非损失函数问题，非特征工程问题\n"
    "→ 论文中可作为探测器角分辨极限的定量结论"
), left=Inches(0.8), top=Inches(1.5), width=Inches(11.7), height=Inches(5.5),
   font_size=Pt(18))
add_slide_number(s, slide_num)


# ══════════════════════════════════════════════════════════════
#  SLIDE 11: Experiment Plan
# ══════════════════════════════════════════════════════════════
slide_num += 1
s = add_slide()
add_title_bar(s, "九、下一步实验计划")

# Dataset info
add_body_text(s, (
    "数据集：9 个角度类别 — 10°, 20°, 30°, 45°, 50°, 60°, 70°, 80°, 90°"
), left=Inches(0.8), top=Inches(1.4), width=Inches(11.7), height=Inches(0.5),
   font_size=Pt(20))

# Ablation matrix table
add_table(s, [
    ["", "交叉熵 + One-hot", "Wasserstein + 软标签", "MSE 回归"],
    ["ResNet-18", "A: 基线", "B: 改损失函数", "E: 回归基线"],
    ["ShallowResNet", "C: 改网络", "D: 两者都改", "F: 回归+浅层"],
], left=Inches(1.5), top=Inches(2.2),
   col_widths=[Inches(2.0), Inches(2.8), Inches(3.2), Inches(2.8)],
   font_size=Pt(16))

add_body_text(s, (
    "评估指标：\n"
    "• 分类方案：准确率、Macro-F1、混淆矩阵\n"
    "• 回归方案：MAE（平均绝对角度误差）、RMSE、R²\n"
    "• 统一对比指标：MAE（度）— 所有方案可比\n"
    "\n"
    "目的：2×3 消融实验，量化每项改进的独立贡献"
), left=Inches(0.8), top=Inches(4.3), width=Inches(11.7), height=Inches(2.8),
   font_size=Pt(18))
add_slide_number(s, slide_num)


# ══════════════════════════════════════════════════════════════
#  SLIDE 12: Summary
# ══════════════════════════════════════════════════════════════
slide_num += 1
s = add_slide()
add_title_bar(s, "十、总结")
add_bullet_slide(s, "", [
    ("1. 论文核心贡献", 0),
    ("仅用 4 像素 + 铅填充 + 神经网络实现 γ 射线 ~1° 方向精度", 1),
    ("关键：物理先验注入（Filter Layer）+ Wasserstein 分布回归", 1),
    ("", 0),
    ("2. 与本课题的关键差异", 0),
    ("论文：方向信息由宏观遮挡人为创造 → 信号对比度高", 1),
    ("本课题：方向信息来自微观径迹自然形态 → 近垂直区间信号极弱", 1),
    ("", 0),
    ("3. 受启发的三项改进", 0),
    ("① Wasserstein 损失 — 让模型学到角度有序性", 1),
    ("② 高斯软标签 — 缓解相邻角度的梯度矛盾", 1),
    ("③ ShallowResNet — 适配稀疏数据，减少信息丢失", 1),
    ("", 0),
    ("4. 下一步", 0),
    ("9 类全角度 2×3 消融实验 → 量化各改进贡献 → 确定最优方案", 1),
], start_top=Inches(1.3))
add_slide_number(s, slide_num)


# ══════════════════════════════════════════════════════════════
#  Save
# ══════════════════════════════════════════════════════════════
prs.save(PPT_PATH)
print(f"Presentation saved to: {PPT_PATH}")
print(f"Total slides: {len(prs.slides)}")
